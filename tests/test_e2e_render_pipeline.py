"""End-to-end render pipeline tests.

Each test exercises a realistic CEO-use-case spec through the full
spec → render → file → QA check pipeline and asserts that:
  - the output file is written and valid
  - the qa_report passes all checks
  - key content (titles, section labels, slide text) survives into the file
"""
from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

import pytest

from src.presentation import DeckSlideSpec, DeckSpec, MemoSectionSpec, MemoSpec
from src.presentation.render_docx import render_docx_memo
from src.presentation.render_pptx import render_pptx_deck
from src.presentation.render_xlsx import render_xlsx_workbook
from src.workflows.workbook_models import (
    WorkbookChartSpec,
    WorkbookSheetSpec,
    WorkbookSpec,
    WorkbookTable,
)


# ---------------------------------------------------------------------------
# xlsx — budget variance workbook
# ---------------------------------------------------------------------------

@pytest.fixture()
def budget_variance_spec() -> WorkbookSpec:
    return WorkbookSpec(
        workbook_title="Q2 Budget Variance Review",
        metadata={"template_id": "finance_workbook_v1", "theme_id": "board_formal"},
        sheets=[
            WorkbookSheetSpec(
                name="Summary",
                kind="summary",
                tables=[
                    WorkbookTable(
                        title="Executive Summary",
                        columns=["Section", "Detail"],
                        rows=[
                            ["Overall Status", "On track — cloud costs 6% favorable vs budget."],
                            ["Top Risk", "Headcount ramp may accelerate burn in Q3."],
                        ],
                    )
                ],
            ),
            WorkbookSheetSpec(
                name="Model",
                kind="model",
                tables=[
                    WorkbookTable(
                        title="Budget vs Actual",
                        columns=["Metric", "Budget", "Actual", "Variance"],
                        rows=[
                            ["Revenue", "4.2M", "4.5M", "300K"],
                            ["Cloud Cost", "520K", "488K", "-32K"],
                            ["Headcount Cost", "1.8M", "1.82M", "20K"],
                            ["Total OpEx", "2.6M", "2.62M", "20K"],
                        ],
                    )
                ],
                chart_specs=[
                    WorkbookChartSpec(
                        title="Budget vs Actual by Category",
                        chart_type="bar",
                        x_axis="Metric",
                        y_axis="Actual",
                        series_label="Actual",
                        source_sheet="Model",
                        source_table="Budget vs Actual",
                    )
                ],
            ),
            WorkbookSheetSpec(
                name="Variance",
                kind="variance",
                tables=[
                    WorkbookTable(
                        title="Period Comparison",
                        columns=["Metric", "Prior Quarter", "Current Quarter"],
                        rows=[
                            ["Revenue", "Q1 2026", "Q2 2026"],
                            ["Burn Rate", "Q1 2026", "Q2 2026"],
                        ],
                    )
                ],
            ),
        ],
    )


def test_xlsx_budget_variance_renders_and_passes_qa(tmp_path: Path, budget_variance_spec: WorkbookSpec) -> None:
    path = tmp_path / "budget_variance.xlsx"
    result = render_xlsx_workbook(path=path, spec=budget_variance_spec, artifact_id="test:1:budget_variance")

    assert path.exists(), "xlsx file was not written"
    assert result["qa_report"]["passed"], (
        f"QA failed: {result['qa_report']['warnings']}"
    )
    assert result["qa_report"]["format"] == "xlsx"

    from openpyxl import load_workbook
    wb = load_workbook(path)
    assert set(wb.sheetnames) == {"Summary", "Model", "Variance"}
    assert wb["Model"]["A1"].value == "Q2 Budget Variance Review"
    assert len(wb["Model"]._charts) == 1


def test_xlsx_qa_report_includes_sheet_checks(tmp_path: Path, budget_variance_spec: WorkbookSpec) -> None:
    path = tmp_path / "budget_variance_checks.xlsx"
    result = render_xlsx_workbook(path=path, spec=budget_variance_spec, artifact_id="test:2:checks")

    check_names = {c["name"] for c in result["qa_report"]["checks"]}
    assert "sheet_present_Summary" in check_names
    assert "sheet_present_Model" in check_names
    assert "sheet_present_Variance" in check_names
    assert "valid_zip" in check_names
    assert all(c["passed"] for c in result["qa_report"]["checks"])


def test_xlsx_period_metadata_extracted(tmp_path: Path, budget_variance_spec: WorkbookSpec) -> None:
    path = tmp_path / "budget_variance_meta.xlsx"
    result = render_xlsx_workbook(path=path, spec=budget_variance_spec, artifact_id="test:3:meta")

    pairs = result["preview_metadata"]["period_coverage"]["comparison_pairs"]
    assert {"prior": "Q1 2026", "current": "Q2 2026"} in pairs


# ---------------------------------------------------------------------------
# docx — board executive memo
# ---------------------------------------------------------------------------

@pytest.fixture()
def board_memo_spec() -> MemoSpec:
    return MemoSpec(
        title="Q2 Board Executive Memo",
        summary="Revenue exceeded plan by 7%. Cloud cost favorable. Headcount risk emerging in Q3.",
        section_order=["Key Findings", "Financial Highlights", "Recommended Actions", "Open Questions"],
        sections=[
            MemoSectionSpec(
                label="Key Findings",
                items=[
                    "Revenue: $4.5M vs $4.2M plan — 7% favorable.",
                    "Cloud cost: $488K vs $520K plan — 6% favorable.",
                    "Headcount: on track, but Q3 ramp may accelerate burn.",
                ],
            ),
            MemoSectionSpec(
                label="Financial Highlights",
                content="Gross margin improved 2pp quarter-over-quarter driven by infrastructure efficiency gains.",
            ),
            MemoSectionSpec(
                label="Recommended Actions",
                items=[
                    "Lock Q3 headcount plan with finance by April 4.",
                    "Approve cloud cost efficiency initiative before next board cycle.",
                    "Prepare scenario analysis for board on Q3 burn sensitivity.",
                ],
            ),
            MemoSectionSpec(
                label="Open Questions",
                items=[
                    "Does the board want a revised FY26 forecast before the next meeting?",
                    "Has legal reviewed the new vendor contract amendment?",
                ],
            ),
        ],
        assumptions=["Board pack reflects the March 31 close."],
        open_questions=["Final legal sign-off pending on vendor amendment."],
        metadata={"template_id": "board_memo_v1", "theme_id": "board_formal"},
    )


def test_docx_board_memo_renders_and_passes_qa(tmp_path: Path, board_memo_spec: MemoSpec) -> None:
    path = tmp_path / "board_memo.docx"
    result = render_docx_memo(path=path, memo_spec=board_memo_spec)

    assert path.exists(), "docx file was not written"
    assert result["qa_report"]["passed"], (
        f"QA failed: {result['qa_report']['warnings']}"
    )
    assert result["qa_report"]["format"] == "docx"


def test_docx_board_memo_content_in_document_xml(tmp_path: Path, board_memo_spec: MemoSpec) -> None:
    path = tmp_path / "board_memo_content.docx"
    render_docx_memo(path=path, memo_spec=board_memo_spec)

    with ZipFile(path) as zf:
        doc_xml = zf.read("word/document.xml").decode("utf-8")

    assert "Q2 Board Executive Memo" in doc_xml
    assert "Key Findings" in doc_xml
    assert "Financial Highlights" in doc_xml
    assert "Recommended Actions" in doc_xml
    assert "Lock Q3 headcount plan" in doc_xml


def test_docx_board_memo_section_order_respected(tmp_path: Path, board_memo_spec: MemoSpec) -> None:
    path = tmp_path / "board_memo_order.docx"
    render_docx_memo(path=path, memo_spec=board_memo_spec)

    with ZipFile(path) as zf:
        doc_xml = zf.read("word/document.xml").decode("utf-8")

    assert doc_xml.index("Key Findings") < doc_xml.index("Financial Highlights")
    assert doc_xml.index("Financial Highlights") < doc_xml.index("Recommended Actions")
    assert doc_xml.index("Recommended Actions") < doc_xml.index("Open Questions")


def test_docx_qa_checks_include_section_labels(tmp_path: Path, board_memo_spec: MemoSpec) -> None:
    path = tmp_path / "board_memo_qa.docx"
    result = render_docx_memo(path=path, memo_spec=board_memo_spec)

    check_names = {c["name"] for c in result["qa_report"]["checks"]}
    assert "section_Key Findings" in check_names
    assert "section_Recommended Actions" in check_names
    assert "title_in_document" in check_names
    assert "has_body_element" in check_names


# ---------------------------------------------------------------------------
# pptx — board deck with mixed slide kinds
# ---------------------------------------------------------------------------

@pytest.fixture()
def board_deck_spec() -> DeckSpec:
    return DeckSpec(
        title="Q2 Board Review Deck",
        slide_order=["Title", "Key Metrics", "Budget vs Actual", "Decision Points", "Appendix"],
        slides=[
            DeckSlideSpec(
                title="Title",
                bullets=["Q2 2026 Board Financial Review", "Prepared by agenticMIND"],
                kind="title",
            ),
            DeckSlideSpec(
                title="Key Metrics",
                bullets=[
                    "Revenue: $4.5M (+7% vs plan)",
                    "Cloud cost: $488K (6% favorable)",
                    "Runway: 18 months",
                ],
                kind="metric",
            ),
            DeckSlideSpec(
                title="Budget vs Actual",
                bullets=[
                    "OpEx on plan at $2.62M vs $2.6M budget.",
                    "No material unfavorable variances above 10% threshold.",
                ],
                kind="content",
            ),
            DeckSlideSpec(
                title="Decision Points",
                bullets=[
                    "Approve Q3 headcount plan (lock by April 4).",
                    "Authorize cloud efficiency initiative ($120K projected savings).",
                ],
                kind="decision",
            ),
            DeckSlideSpec(
                title="Appendix",
                bullets=["Detailed metric support tables available on request."],
                kind="appendix",
            ),
        ],
        metadata={"template_id": "board_deck_v1", "theme_id": "board_formal"},
    )


def test_pptx_board_deck_renders_and_passes_qa(tmp_path: Path, board_deck_spec: DeckSpec) -> None:
    path = tmp_path / "board_deck.pptx"
    result = render_pptx_deck(path=path, deck_spec=board_deck_spec)

    assert path.exists(), "pptx file was not written"
    assert result["qa_report"]["passed"], (
        f"QA failed: {result['qa_report']['warnings']}"
    )
    assert result["qa_report"]["format"] == "pptx"
    assert result["preview_metadata"]["slide_count"] == 5


def test_pptx_board_deck_slide_order_and_content(tmp_path: Path, board_deck_spec: DeckSpec) -> None:
    path = tmp_path / "board_deck_order.pptx"
    render_pptx_deck(path=path, deck_spec=board_deck_spec)

    from pptx import Presentation

    prs = Presentation(str(path))
    slide_text = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        slide_text.append("\n".join(parts))

    assert "Title" in slide_text[0]
    assert "Key Metrics" in slide_text[1]
    assert "Decision Points" in slide_text[3]
    assert "Appendix" in slide_text[4]


def test_pptx_board_deck_qa_checks_all_slides(tmp_path: Path, board_deck_spec: DeckSpec) -> None:
    path = tmp_path / "board_deck_qa.pptx"
    result = render_pptx_deck(path=path, deck_spec=board_deck_spec)

    check_names = {c["name"] for c in result["qa_report"]["checks"]}
    for i in range(1, 6):
        assert f"slide_{i}_title" in check_names, f"missing slide_{i}_title check"

    assert "slide_count" in check_names
    assert "valid_zip" in check_names
    assert "member_presentation_xml" in check_names


def test_pptx_operator_theme_renders_and_passes_qa(tmp_path: Path, board_deck_spec: DeckSpec) -> None:
    path = tmp_path / "operator_deck.pptx"
    operator_spec = board_deck_spec.model_copy(
        update={"metadata": {"template_id": "board_deck_v1", "theme_id": "operator_modern"}}
    )
    result = render_pptx_deck(path=path, deck_spec=operator_spec)

    assert result["qa_report"]["passed"], f"QA failed: {result['qa_report']['warnings']}"
    check_names = {c["name"] for c in result["qa_report"]["checks"]}
    assert "python_pptx_openable" in check_names


# ---------------------------------------------------------------------------
# cross-format: qa_report structure contract
# ---------------------------------------------------------------------------

def test_qa_report_structure_contract(tmp_path: Path) -> None:
    """qa_report dict must always contain the required keys across all three formats."""
    xlsx_path = tmp_path / "mini.xlsx"
    docx_path = tmp_path / "mini.docx"
    pptx_path = tmp_path / "mini.pptx"

    xlsx_result = render_xlsx_workbook(
        path=xlsx_path,
        spec=WorkbookSpec(
            workbook_title="Mini",
            metadata={"template_id": "finance_workbook_v1"},
            sheets=[WorkbookSheetSpec(name="Summary", kind="summary")],
        ),
        artifact_id="test:contract:xlsx",
    )
    docx_result = render_docx_memo(
        path=docx_path,
        memo_spec=MemoSpec(
            title="Mini Memo",
            summary="Brief.",
            sections=[MemoSectionSpec(label="Findings", items=["All good."])],
            metadata={"template_id": "board_memo_v1", "theme_id": "board_formal"},
        ),
    )
    pptx_result = render_pptx_deck(
        path=pptx_path,
        deck_spec=DeckSpec(
            title="Mini Deck",
            slides=[DeckSlideSpec(title="Overview", bullets=["One slide."])],
            metadata={"template_id": "board_deck_v1", "theme_id": "board_formal"},
        ),
    )

    for fmt, result in (("xlsx", xlsx_result), ("docx", docx_result), ("pptx", pptx_result)):
        report = result["qa_report"]
        assert "passed" in report, f"{fmt}: missing 'passed'"
        assert "format" in report, f"{fmt}: missing 'format'"
        assert "checks" in report, f"{fmt}: missing 'checks'"
        assert "warnings" in report, f"{fmt}: missing 'warnings'"
        assert isinstance(report["checks"], list), f"{fmt}: 'checks' must be a list"
        assert report["format"] == fmt, f"{fmt}: format mismatch"
        assert report["passed"] is True, f"{fmt}: expected passing QA, got: {report['warnings']}"
