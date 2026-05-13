import logging
from pathlib import Path

from pptx import Presentation

from src.presentation import DeckSlideSpec, DeckSpec
from src.presentation.render_pptx import render_pptx_deck
from src.tools.base import ToolContext
from src.tools.document_tools import CreatePptxDeckTool


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    return "\n".join(parts)


def test_render_pptx_deck_writes_file_and_preview(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    spec = DeckSpec(
        title="Meeting Prep Deck",
        slide_order=["Context", "Recommended Actions"],
        slides=[
            DeckSlideSpec(title="Recommended Actions", bullets=["Confirm hiring plan."]),
            DeckSlideSpec(title="Context", bullets=["Quarterly board meeting next Tuesday."]),
        ],
        metadata={"template_id": "meeting_prep_deck_v1", "theme_id": "board_formal"},
    )

    result = render_pptx_deck(path=path, deck_spec=spec)
    prs = Presentation(str(path))

    assert path.exists()
    assert result["preview_format"] == "md"
    assert result["preview_metadata"]["slide_count"] == 2
    assert len(prs.slides) == 2
    assert "Context" in _slide_text(prs.slides[0])
    assert "Recommended Actions" in _slide_text(prs.slides[1])


def test_render_pptx_deck_generates_openable_powerpoint_with_slide_kinds(tmp_path: Path) -> None:
    path = tmp_path / "board-deck.pptx"
    spec = DeckSpec(
        title="Board Deck",
        slides=[
            DeckSlideSpec(title="Title", bullets=["Q2 board review"], kind="title"),
            DeckSlideSpec(title="Key Metrics", bullets=["Revenue up 18%"], kind="metric"),
            DeckSlideSpec(title="Decision Points", bullets=["Approve hiring plan"], kind="decision"),
            DeckSlideSpec(title="Appendix", bullets=["Metric support"], kind="appendix"),
        ],
        metadata={"template_id": "board_deck_v1", "theme_id": "board_formal"},
    )

    result = render_pptx_deck(path=path, deck_spec=spec)
    prs = Presentation(str(path))

    assert result["qa_report"]["passed"] is True
    assert len(prs.slides) == 4
    assert "Q2 board review" in _slide_text(prs.slides[0])
    assert "Metric Summary" in _slide_text(prs.slides[1])
    assert "Decision Summary" in _slide_text(prs.slides[2])
    assert "Appendix" in _slide_text(prs.slides[3])


def test_render_pptx_cover_style_changes_theme_output(tmp_path: Path) -> None:
    board_path = tmp_path / "board-style.pptx"
    operator_path = tmp_path / "operator-style.pptx"
    board_spec = DeckSpec(
        title="Board Deck",
        slides=[
            DeckSlideSpec(title="Title", bullets=["Q2 board review"], kind="title"),
            DeckSlideSpec(title="Key Metrics", bullets=["Revenue up 18%"], kind="metric"),
        ],
        metadata={"template_id": "board_deck_v1", "theme_id": "board_formal"},
    )
    operator_spec = DeckSpec(
        title="Operator Deck",
        slides=[
            DeckSlideSpec(title="Title", bullets=["Weekly operating review"], kind="title"),
            DeckSlideSpec(title="Key Metrics", bullets=["Conversion up 6%"], kind="metric"),
        ],
        metadata={"template_id": "board_deck_v1", "theme_id": "operator_modern"},
    )

    render_pptx_deck(path=board_path, deck_spec=board_spec)
    render_pptx_deck(path=operator_path, deck_spec=operator_spec)

    board_prs = Presentation(str(board_path))
    operator_prs = Presentation(str(operator_path))

    assert "Q2 board review" in _slide_text(board_prs.slides[0])
    assert "Weekly operating review" in _slide_text(operator_prs.slides[0])
    assert _slide_text(board_prs.slides[1]) != _slide_text(operator_prs.slides[1])


def test_create_pptx_deck_tool_delegates_to_renderer(tmp_path: Path) -> None:
    tool = CreatePptxDeckTool()
    context = ToolContext(interaction_id=55, ceo_id="ceo_test")
    result = tool.invoke(
        context,
        output_path=str(tmp_path / "tool-deck.pptx"),
        deck_spec={
            "title": "Meeting Prep Deck",
            "slides": [{"title": "Context", "bullets": ["Board prep is underway."]}],
            "metadata": {"template_id": "meeting_prep_deck_v1", "theme_id": "board_formal"},
        },
    )

    assert result.success is True
    assert Path(result.data["path"]).exists()
    assert result.metadata["preview_metadata"]["slide_count"] == 1
    assert len(Presentation(str(result.data["path"])).slides) == 1


def test_create_pptx_deck_tool_accepts_top_level_fields(tmp_path: Path) -> None:
    tool = CreatePptxDeckTool()
    context = ToolContext(interaction_id=56, ceo_id="ceo_test")
    result = tool.invoke(
        context,
        output_path=str(tmp_path / "tool-deck-top-level.pptx"),
        title="Meeting Prep Deck",
        slides=[{"title": "Context", "bullets": ["Board prep is underway."]}],
        template_id="meeting_prep_deck_v1",
        theme_id="board_formal",
    )

    assert result.success is True
    assert Path(result.data["path"]).exists()
    assert result.metadata["preview_metadata"]["slide_count"] == 1
    assert len(Presentation(str(result.data["path"])).slides) == 1


def test_create_pptx_deck_logs_invalid_payload_keys(caplog, tmp_path: Path) -> None:
    tool = CreatePptxDeckTool()
    context = ToolContext(interaction_id=57, ceo_id="ceo_test", stage="synthesizer")

    with caplog.at_level(logging.WARNING):
        result = tool.invoke(
            context,
            output_path=str(tmp_path / "invalid-deck.pptx"),
            deck_spec={"slides": [{"title": "Context"}]},
        )

    assert result.success is False
    assert "Invalid deck payload" in (result.error or "")
    assert "Invalid deck payload for create_pptx_deck" in caplog.text
    assert "payload_keys=['slides']" in caplog.text
    assert "kwargs_keys=['deck_spec', 'output_path']" in caplog.text
