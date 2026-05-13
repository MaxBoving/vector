from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.presentation import DeckSpec, deck_spec_to_preview_markdown, get_artifact_template, resolve_brand_theme
from src.presentation.artifact_contracts import DEFAULT_DECK_TEMPLATE_ID, normalize_deck_spec
from src.presentation.render_qa import qa_check_pptx


SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def _ordered_slides(deck_spec: DeckSpec) -> list:
    template = get_artifact_template(str(deck_spec.metadata.get("template_id") or DEFAULT_DECK_TEMPLATE_ID))
    ordered_titles = list(deck_spec.slide_order) or list(getattr(template, "slide_sequence", []))
    seen = set(ordered_titles)
    ordered_titles.extend([slide.title for slide in deck_spec.slides if slide.title not in seen])
    slide_map = {slide.title: slide for slide in deck_spec.slides}
    ordered = [slide_map[title] for title in ordered_titles if title in slide_map]
    return ordered or list(deck_spec.slides)


def _font_size_points(cover_style: str) -> dict[str, float]:
    if cover_style == "formal":
        return {
            "title": 28,
            "subtitle": 15,
            "kicker": 11,
            "section_title": 23,
            "appendix_title": 17,
            "body": 12,
        }
    if cover_style == "operator":
        return {
            "title": 25,
            "subtitle": 13.5,
            "kicker": 10,
            "section_title": 21,
            "appendix_title": 16,
            "body": 11.5,
        }
    return {
        "title": 26,
        "subtitle": 14,
        "kicker": 10.5,
        "section_title": 22,
        "appendix_title": 16.5,
        "body": 11.8,
    }


def _surface_tokens(theme) -> dict[str, RGBColor]:
    if theme.cover_style == "formal":
        return {
            "title_background": _rgb(theme.colors.accent),
            "content_background": _rgb(theme.colors.background),
            "content_panel": _rgb(theme.colors.accent),
            "appendix_panel": _rgb(theme.colors.surface),
            "appendix_border": _rgb(theme.colors.border),
        }
    if theme.cover_style == "operator":
        return {
            "title_background": _rgb(theme.colors.background),
            "content_background": _rgb(theme.colors.surface),
            "content_panel": _rgb(theme.colors.surface),
            "appendix_panel": _rgb(theme.colors.background),
            "appendix_border": _rgb(theme.colors.accent),
        }
    return {
        "title_background": _rgb(theme.colors.accent),
        "content_background": _rgb(theme.colors.surface),
        "content_panel": _rgb(theme.colors.accent),
        "appendix_panel": _rgb(theme.colors.surface),
        "appendix_border": _rgb(theme.colors.border),
    }


def _set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    *,
    left,
    top,
    width,
    height,
    text: str,
    font_name: str,
    font_size: float,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical_anchor
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_panel(slide, *, left, top, width, height, fill_color: RGBColor, line_color: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.adjustments[0] = 0.08
    return shape


def _populate_body_text(shape, bullets: list[str], *, font_name: str, font_size: float, color: RGBColor, bold: bool = False, level: int = 0) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = level
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(6)
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def _add_title_slide(prs: Presentation, slide_spec, *, theme, sizes: dict[str, float], surfaces: dict[str, RGBColor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, surfaces["title_background"])

    title_color = _rgb(theme.colors.primary)
    body_color = _rgb(theme.colors.text)

    _add_textbox(
        slide,
        left=Inches(0.8),
        top=Inches(1.25),
        width=Inches(8.4),
        height=Inches(1.5),
        text=slide_spec.title,
        font_name=theme.typography.heading_family,
        font_size=sizes["title"],
        color=title_color,
        bold=theme.typography.title_weight >= 600,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    subtitle = slide_spec.bullets[0] if slide_spec.bullets else ""
    if subtitle:
        _add_textbox(
            slide,
            left=Inches(0.85),
            top=Inches(3.15),
            width=Inches(8.1),
            height=Inches(1.6),
            text=subtitle,
            font_name=theme.typography.body_family,
            font_size=sizes["subtitle"],
            color=body_color,
            bold=theme.typography.body_weight >= 600,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )


def _add_content_slide(prs: Presentation, slide_spec, *, theme, sizes: dict[str, float], surfaces: dict[str, RGBColor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, _rgb(theme.colors.background))

    _add_textbox(
        slide,
        left=Inches(0.75),
        top=Inches(0.7),
        width=Inches(8.4),
        height=Inches(0.8),
        text=slide_spec.title,
        font_name=theme.typography.heading_family,
        font_size=sizes["section_title"],
        color=_rgb(theme.colors.primary),
        bold=theme.typography.heading_weight >= 600,
    )

    body_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.65), Inches(8.25), Inches(4.95))
    _populate_body_text(
        body_box,
        slide_spec.bullets,
        font_name=theme.typography.body_family,
        font_size=sizes["body"],
        color=_rgb(theme.colors.text),
        bold=theme.typography.body_weight >= 600,
    )


def _add_metric_or_decision_slide(prs: Presentation, slide_spec, *, theme, sizes: dict[str, float], surfaces: dict[str, RGBColor], kind: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, surfaces["content_background"])

    kicker = "Metric Summary" if kind == "metric" else "Decision Summary"
    _add_textbox(
        slide,
        left=Inches(0.82),
        top=Inches(0.6),
        width=Inches(3.0),
        height=Inches(0.45),
        text=kicker,
        font_name=theme.typography.heading_family,
        font_size=sizes["kicker"],
        color=_rgb(theme.colors.secondary),
        bold=theme.typography.heading_weight >= 600,
    )

    _add_textbox(
        slide,
        left=Inches(0.8),
        top=Inches(1.0),
        width=Inches(8.3),
        height=Inches(0.8),
        text=slide_spec.title,
        font_name=theme.typography.heading_family,
        font_size=sizes["section_title"],
        color=_rgb(theme.colors.primary),
        bold=theme.typography.heading_weight >= 600,
    )

    panel = _add_panel(
        slide,
        left=Inches(0.8),
        top=Inches(1.9),
        width=Inches(8.35),
        height=Inches(4.95),
        fill_color=surfaces["content_panel"],
        line_color=surfaces["content_panel"],
    )
    _populate_body_text(
        panel,
        slide_spec.bullets,
        font_name=theme.typography.body_family,
        font_size=sizes["body"],
        color=_rgb(theme.colors.text),
        bold=theme.typography.body_weight >= 600,
    )


def _add_appendix_slide(prs: Presentation, slide_spec, *, theme, sizes: dict[str, float], surfaces: dict[str, RGBColor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, surfaces["appendix_panel"])

    _add_textbox(
        slide,
        left=Inches(0.82),
        top=Inches(0.72),
        width=Inches(8.25),
        height=Inches(0.7),
        text=slide_spec.title,
        font_name=theme.typography.heading_family,
        font_size=sizes["appendix_title"],
        color=_rgb(theme.colors.secondary),
        bold=theme.typography.heading_weight >= 600,
    )

    panel = _add_panel(
        slide,
        left=Inches(0.8),
        top=Inches(1.55),
        width=Inches(8.35),
        height=Inches(5.1),
        fill_color=surfaces["appendix_panel"],
        line_color=surfaces["appendix_border"],
    )
    _populate_body_text(
        panel,
        slide_spec.bullets,
        font_name=theme.typography.body_family,
        font_size=sizes["body"],
        color=_rgb(theme.colors.text),
        bold=theme.typography.body_weight >= 600,
    )


def _render_slide(prs: Presentation, slide_spec, *, theme, sizes: dict[str, float], surfaces: dict[str, RGBColor]) -> None:
    if slide_spec.kind == "title":
        _add_title_slide(prs, slide_spec, theme=theme, sizes=sizes, surfaces=surfaces)
        return
    if slide_spec.kind in {"metric", "decision"}:
        _add_metric_or_decision_slide(
            prs,
            slide_spec,
            theme=theme,
            sizes=sizes,
            surfaces=surfaces,
            kind=slide_spec.kind,
        )
        return
    if slide_spec.kind == "appendix":
        _add_appendix_slide(prs, slide_spec, theme=theme, sizes=sizes, surfaces=surfaces)
        return
    _add_content_slide(prs, slide_spec, theme=theme, sizes=sizes, surfaces=surfaces)


def render_pptx_deck(*, path: Path, deck_spec: DeckSpec) -> dict[str, object]:
    deck_spec = normalize_deck_spec(deck_spec)
    theme = resolve_brand_theme(deck_spec.metadata.get("theme_id"))
    ordered_slides = _ordered_slides(deck_spec)
    sizes = _font_size_points(theme.cover_style)
    surfaces = _surface_tokens(theme)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    while prs.slides:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for slide_spec in ordered_slides:
        _render_slide(prs, slide_spec, theme=theme, sizes=sizes, surfaces=surfaces)

    prs.core_properties.title = deck_spec.title
    prs.core_properties.author = "agenticMIND"
    prs.save(path)

    qa_report = qa_check_pptx(path, deck_spec)
    return {
        "preview_content": deck_spec_to_preview_markdown(deck_spec),
        "preview_format": "md",
        "preview_metadata": {
            **deck_spec.metadata,
            "slide_count": len(ordered_slides),
        },
        "qa_report": qa_report.model_dump(),
    }
