from __future__ import annotations

from pathlib import Path
from typing import Any
from zipfile import BadZipFile, ZipFile

from pydantic import BaseModel, Field


class QACheckResult(BaseModel):
    name: str
    passed: bool
    detail: str = ""


class RenderQAReport(BaseModel):
    passed: bool
    format: str
    checks: list[QACheckResult] = Field(default_factory=list)
    warnings: list[str] = Field(default_factory=list)

    @property
    def failed_checks(self) -> list[QACheckResult]:
        return [c for c in self.checks if not c.passed]


def _check(name: str, condition: bool, detail: str = "") -> QACheckResult:
    return QACheckResult(name=name, passed=condition, detail=detail)


def _safe_load_zip(path: Path) -> tuple[list[str], str | None]:
    """Return (namelist, error). error is None on success."""
    try:
        with ZipFile(path) as zf:
            return zf.namelist(), None
    except BadZipFile as exc:
        return [], str(exc)


# ---------------------------------------------------------------------------
# xlsx
# ---------------------------------------------------------------------------

def qa_check_xlsx(path: Path, spec: Any) -> RenderQAReport:
    checks: list[QACheckResult] = []
    warnings: list[str] = []

    checks.append(_check("file_exists", path.exists(), str(path)))
    if not path.exists():
        return RenderQAReport(passed=False, format="xlsx", checks=checks, warnings=warnings)

    file_size = path.stat().st_size
    checks.append(_check("file_nonempty", file_size > 0, f"{file_size} bytes"))

    names, zip_error = _safe_load_zip(path)
    checks.append(_check("valid_zip", zip_error is None, zip_error or f"{len(names)} members"))
    if zip_error:
        return RenderQAReport(passed=False, format="xlsx", checks=checks, warnings=warnings)

    for member in ("xl/workbook.xml", "[Content_Types].xml"):
        safe_name = member.replace("/", "_").replace(".", "_").replace("[", "").replace("]", "")
        checks.append(_check(f"member_{safe_name}", member in names, member))

    expected_sheets = spec.sheets or []
    if expected_sheets:
        try:
            from openpyxl import load_workbook as _load_wb

            wb = _load_wb(path)
            actual_count = len(wb.sheetnames)
            checks.append(_check(
                "sheet_count",
                actual_count == len(expected_sheets),
                f"expected {len(expected_sheets)}, got {actual_count}",
            ))
            for sheet in expected_sheets:
                sheet_name = str(sheet.name or "")[:31]
                checks.append(_check(f"sheet_present_{sheet_name}", sheet_name in wb.sheetnames, sheet_name))

            for sname in wb.sheetnames:
                ws = wb[sname]
                checks.append(_check(f"sheet_a1_content_{sname}", ws["A1"].value is not None, sname))
        except Exception as exc:
            warnings.append(f"Could not load workbook for detailed checks: {exc}")
            checks.append(_check("workbook_readable", False, str(exc)))

    passed = all(c.passed for c in checks)
    warnings.extend(f"FAILED: {c.name} — {c.detail}" for c in checks if not c.passed)
    return RenderQAReport(passed=passed, format="xlsx", checks=checks, warnings=warnings)


# ---------------------------------------------------------------------------
# docx
# ---------------------------------------------------------------------------

def qa_check_docx(path: Path, spec: Any) -> RenderQAReport:
    checks: list[QACheckResult] = []
    warnings: list[str] = []

    checks.append(_check("file_exists", path.exists(), str(path)))
    if not path.exists():
        return RenderQAReport(passed=False, format="docx", checks=checks, warnings=warnings)

    file_size = path.stat().st_size
    checks.append(_check("file_nonempty", file_size > 0, f"{file_size} bytes"))

    names, zip_error = _safe_load_zip(path)
    checks.append(_check("valid_zip", zip_error is None, zip_error or f"{len(names)} members"))
    if zip_error:
        return RenderQAReport(passed=False, format="docx", checks=checks, warnings=warnings)

    checks.append(_check("member_word_document_xml", "word/document.xml" in names))
    checks.append(_check("member_styles_xml", "word/styles.xml" in names))
    checks.append(_check("member_document_rels", "word/_rels/document.xml.rels" in names))
    checks.append(_check("member_theme_xml", "word/theme/theme1.xml" in names))
    checks.append(_check("member_settings_xml", "word/settings.xml" in names))

    try:
        with ZipFile(path) as zf:
            document_xml = zf.read("word/document.xml").decode("utf-8")
    except KeyError:
        warnings.append("word/document.xml not readable")
        return RenderQAReport(passed=False, format="docx", checks=checks, warnings=warnings)

    checks.append(_check("has_body_element", "<w:body>" in document_xml))

    title = str(spec.title or "")
    if title:
        checks.append(_check("title_in_document", title in document_xml, f"'{title}'"))

    summary = str(spec.summary or "")
    if summary:
        snippet = summary[:60]
        checks.append(_check("summary_in_document", snippet in document_xml, f"'{snippet}'"))

    for section in spec.sections or []:
        label = str(section.label or "")
        if label:
            checks.append(_check(f"section_{label[:24]}", label in document_xml, label))

    passed = all(c.passed for c in checks)
    warnings.extend(f"FAILED: {c.name} — {c.detail}" for c in checks if not c.passed)
    return RenderQAReport(passed=passed, format="docx", checks=checks, warnings=warnings)


# ---------------------------------------------------------------------------
# pptx
# ---------------------------------------------------------------------------

def qa_check_pptx(path: Path, spec: Any) -> RenderQAReport:
    checks: list[QACheckResult] = []
    warnings: list[str] = []

    checks.append(_check("file_exists", path.exists(), str(path)))
    if not path.exists():
        return RenderQAReport(passed=False, format="pptx", checks=checks, warnings=warnings)

    file_size = path.stat().st_size
    checks.append(_check("file_nonempty", file_size > 0, f"{file_size} bytes"))

    names, zip_error = _safe_load_zip(path)
    checks.append(_check("valid_zip", zip_error is None, zip_error or f"{len(names)} members"))
    if zip_error:
        return RenderQAReport(passed=False, format="pptx", checks=checks, warnings=warnings)

    checks.append(_check("member_presentation_xml", "ppt/presentation.xml" in names))
    checks.append(_check("member_content_types", "[Content_Types].xml" in names))

    slides = list(spec.slides or [])
    slide_order = list(spec.slide_order or [])
    if slide_order:
        slide_map = {s.title: s for s in slides}
        ordered = [slide_map[t] for t in slide_order if t in slide_map]
        ordered += [s for s in slides if s.title not in slide_order]
    else:
        ordered = slides

    actual_slide_files = sorted(n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
    checks.append(_check(
        "slide_count",
        len(actual_slide_files) == len(ordered),
        f"expected {len(ordered)}, got {len(actual_slide_files)}",
    ))

    try:
        from pptx import Presentation

        parsed = Presentation(str(path))
        checks.append(_check("python_pptx_openable", True, f"{len(parsed.slides)} slides"))
        checks.append(_check(
            "python_pptx_slide_count",
            len(parsed.slides) == len(ordered),
            f"expected {len(ordered)}, got {len(parsed.slides)}",
        ))
    except Exception as exc:
        checks.append(_check("python_pptx_openable", False, str(exc)))

    with ZipFile(path) as zf:
        for idx, slide in enumerate(ordered):
            slide_path = f"ppt/slides/slide{idx + 1}.xml"
            if slide_path in names:
                slide_xml = zf.read(slide_path).decode("utf-8")
                title = str(slide.title or "")
                if title:
                    checks.append(_check(f"slide_{idx + 1}_title", title in slide_xml, title))
            else:
                checks.append(_check(f"slide_{idx + 1}_exists", False, slide_path))

    passed = all(c.passed for c in checks)
    warnings.extend(f"FAILED: {c.name} — {c.detail}" for c in checks if not c.passed)
    return RenderQAReport(passed=passed, format="pptx", checks=checks, warnings=warnings)
